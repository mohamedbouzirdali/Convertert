import re
import spacy
import pytesseract
import cv2
import numpy as np
import PyPDF2
import docx
from pptx import Presentation

# Optional: Use pdf2image for OCR on PDF pages if text extraction is poor
try:
    from pdf2image import convert_from_bytes
except ImportError:
    convert_from_bytes = None

# Load the spaCy model (make sure to download it via: python -m spacy download en_core_web_sm)
nlp = spacy.load("en_core_web_sm")


def parse_pdf(file):
    """
    Extract text from a PDF. If the built-in text extraction yields little text,
    try converting pages to images and use OCR.
    """
    file.seek(0)
    pdf_reader = PyPDF2.PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        page_text = page.extract_text() or ""
        text += page_text

    # Fallback: if very little text was extracted, try OCR (if pdf2image is available)
    if len(text.strip()) < 50 and convert_from_bytes is not None:
        file.seek(0)
        pages = convert_from_bytes(file.read())
        text = ""
        for page_image in pages:
            # Convert PIL image to an OpenCV image (BGR format)
            open_cv_image = cv2.cvtColor(np.array(page_image), cv2.COLOR_RGB2BGR)
            text += pytesseract.image_to_string(open_cv_image)
    return extract_information(text)


def parse_image(file):
    """
    Use pytesseract OCR on an image file.
    """
    file_bytes = np.asarray(bytearray(file.read()), dtype=np.uint8)
    image = cv2.imdecode(file_bytes, cv2.IMREAD_COLOR)
    text = pytesseract.image_to_string(image)
    return extract_information(text)


def parse_pptx(file):
    """
    Extract text from each slide of a PPTX file.
    """
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + " "
    return extract_information(text)


def parse_docx(file):
    """
    Extract text from a DOCX file.
    """
    document = docx.Document(file)
    text = " ".join([para.text for para in document.paragraphs])
    return extract_information(text)


def extract_information(raw_text):
    """
    Process the raw text using spaCy and custom heuristics to extract structured resume data.
    """
    doc = nlp(raw_text)

    personal_info = extract_personal_info(raw_text, doc)
    education = extract_education(raw_text)
    experience = extract_experience(raw_text)
    skills = extract_skills(raw_text)
    certifications = extract_certifications(raw_text)
    tools = extract_tools(raw_text)
    languages = extract_languages(raw_text)
    summary = extract_summary(raw_text)
    courses_conferences = extract_courses_conferences(raw_text)
    interests = extract_interests(raw_text)

    result = {
        "name": personal_info.get("name", "Unknown"),
        "email": personal_info.get("email", "Unknown"),
        "phone": personal_info.get("phone", "Unknown"),
        "linkedin": personal_info.get("linkedin", "Unknown"),
        "address": personal_info.get("address", "Unknown"),
        "education": education,
        "skills": skills,
        "experience": experience,
        "languages": languages,
        "dob": personal_info.get("dob", "Unknown"),
        "certifications": certifications,
        "tools": tools,
        "summary": summary,
        "interests": interests,
        "courses_conferences": courses_conferences
    }
    return result


def extract_personal_info(text, doc):
    """
    Extracts personal information such as name, email, phone, LinkedIn URL,
    address, and date of birth.
    """
    info = {}

    # Email and phone via regex
    email_match = re.search(r'[\w\.-]+@[\w\.-]+\.\w+', text)
    phone_match = re.search(r'\+?\d[\d\s\-]{7,15}', text)
    info["email"] = email_match.group(0) if email_match else "Unknown"
    info["phone"] = phone_match.group(0) if phone_match else "Unknown"

    # LinkedIn URL via regex
    linkedin_match = re.search(r'(https?://(www\.)?linkedin\.com/in/[\w\-]+)', text, re.IGNORECASE)
    info["linkedin"] = linkedin_match.group(0) if linkedin_match else "Unknown"

    # Use spaCy to extract a person's name (first found PERSON entity)
    for ent in doc.ents:
        if ent.label_ == "PERSON":
            info["name"] = ent.text
            break
    if "name" not in info:
        info["name"] = "Unknown"

    # Simple heuristic for address (e.g., number + street)
    address_match = re.search(r'\d{1,4}\s+\w+\s+(Street|St|Avenue|Ave|Road|Rd)', text, re.IGNORECASE)
    info["address"] = address_match.group(0) if address_match else "Unknown"

    # Try to extract date of birth (if formatted as YYYY-MM-DD)
    dob_match = re.search(r'\b(?:DOB|Date of Birth)[:\s]+(\d{4}-\d{2}-\d{2})', text)
    info["dob"] = dob_match.group(1) if dob_match else "Unknown"

    return info


def extract_education(text):
    """
    Extract education-related details using common degree patterns and university names.
    """
    education_data = []
    education_patterns = [
        r'(Bachelor(?:\'s)? of [\w\s]+)',
        r'(Master(?:\'s)? of [\w\s]+)',
        r'(B\.Sc\. in [\w\s]+)',
        r'(M\.Sc\. in [\w\s]+)',
        r'(Ph\.D\. in [\w\s]+)',
        r'(Associate(?:\'s)? degree in [\w\s]+)'
    ]
    for pattern in education_patterns:
        matches = re.findall(pattern, text, re.IGNORECASE)
        for match in matches:
            if match not in education_data:
                education_data.append(match.strip())

    # Look for university names (a simple heuristic)
    uni_matches = re.findall(r'(University of [\w\s]+)', text, re.IGNORECASE)
    for match in uni_matches:
        if match.strip() not in education_data:
            education_data.append(match.strip())

    return education_data if education_data else "Not specified"


def extract_experience(text):
    """
    Extract experience by looking for sentences containing dates and common job-related keywords.
    """
    experience_data = []
    # Split the text into sentences
    sentences = re.split(r'(?<=[.!?]) +', text)
    for sentence in sentences:
        if re.search(r'\b(20\d{2}|19\d{2})\b', sentence) and any(
            keyword in sentence.lower() for keyword in ["experience", "worked", "intern", "engineer", "manager", "developer"]
        ):
            experience_data.append(sentence.strip())
    return experience_data if experience_data else "Not specified"


def extract_skills(text):
    """
    Extract skills by matching against a pre-defined list and checking for a dedicated 'Skills' section.
    """
    common_skills = [
        "Python", "Java", "C++", "SQL", "Machine Learning", "Deep Learning",
        "NLP", "Data Analysis", "TensorFlow", "Keras", "PyTorch", "Git", "Docker"
    ]
    skills_found = set()
    for skill in common_skills:
        if re.search(r'\b' + re.escape(skill) + r'\b', text, re.IGNORECASE):
            skills_found.add(skill)
    # Look for a comma-separated skills section
    skills_section = re.search(r'Skills[:\-]\s*(.+)', text, re.IGNORECASE)
    if skills_section:
        extra_skills = skills_section.group(1).split(',')
        for s in extra_skills:
            s = s.strip()
            if s:
                skills_found.add(s)
    return list(skills_found) if skills_found else "Not specified"


def extract_certifications(text):
    """
    Extract certification details using simple keyword patterns.
    """
    certs = []
    cert_pattern = r'(Certified [\w\s]+|[A-Z]{2,}\s+Certification)'
    matches = re.findall(cert_pattern, text, re.IGNORECASE)
    for match in matches:
        certs.append(match.strip())
    return certs if certs else "Not specified"


def extract_tools(text):
    """
    Extract tool names by matching against a pre-defined list.
    """
    common_tools = ["PyPDF2", "OpenCV", "TensorFlow", "Keras", "Scikit-learn", "Docker", "Git", "JIRA", "Confluence"]
    tools_found = set()
    for tool in common_tools:
        if re.search(r'\b' + re.escape(tool) + r'\b', text, re.IGNORECASE):
            tools_found.add(tool)
    return list(tools_found) if tools_found else "Not specified"


def extract_languages(text):
    """
    Extract spoken or programming languages mentioned in the resume.
    """
    languages_list = ["English", "Spanish", "French", "German", "Mandarin", "Hindi"]
    languages_found = set()
    for lang in languages_list:
        if re.search(r'\b' + re.escape(lang) + r'\b', text, re.IGNORECASE):
            languages_found.add(lang)
    return list(languages_found) if languages_found else "Not specified"


def extract_summary(text):
    """
    Return the first few sentences as a summary (if available).
    """
    sentences = re.split(r'(?<=[.!?])\s+', text.strip())
    if sentences:
        # Return first 2-3 sentences as summary if they are sufficiently long
        summary = " ".join(sentences[:3])
        return summary if len(summary) > 30 else "Not specified"
    return "Not specified"


def extract_interests(text):
    """
    Extract interests if explicitly mentioned.
    """
    interests_match = re.search(r'Interests[:\-]\s*(.+)', text, re.IGNORECASE)
    if interests_match:
        interests = [i.strip() for i in interests_match.group(1).split(',')]
        return interests
    return "Not specified"


def extract_courses_conferences(text):
    """
    Extract courses or conferences details if available.
    """
    courses_match = re.search(r'(Courses|Conferences)[:\-]\s*(.+)', text, re.IGNORECASE)
    if courses_match:
        courses = [c.strip() for c in courses_match.group(2).split(',')]
        return courses
    return "Not specified"
